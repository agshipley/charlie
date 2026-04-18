"""
Format-aware extraction pipeline for Field Work artifacts.

Supports: .docx, .xlsx, .pdf, .pptx, .md, .txt
Normalized output schema:
    {
        artifact_id: str,
        extracted_at: str (ISO),
        source_format: str,
        title_extracted: str | None,
        sections: [{heading, level, content}],
        tables: [{section_index, rows}],
        full_text: str,
        word_count: int,
        extraction_notes: [str],
    }
"""

import json
import os
import re
from datetime import datetime
from pathlib import Path
from .config import config
from .logging import get_logger

_log = get_logger(__name__)


def extract_artifact(artifact_id: str, original_path: Path, file_extension: str) -> dict:
    """
    Dispatch to the correct format extractor, save the result to
    data/field/extracted/{artifact_id}.json, and return the dict.

    Raises on failure — caller handles the exception and updates artifact status.
    """
    ext = file_extension.lower().lstrip(".")
    dispatch = {
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pdf":  _extract_pdf,
        "pptx": _extract_pptx,
        "md":   _extract_markdown,
        "txt":  _extract_txt,
    }

    _log.info("extraction_started", artifact_id=artifact_id, ext=ext, path=str(original_path))

    extractor = dispatch.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file format: {ext!r}")

    result = extractor(original_path)

    if not result["sections"]:
        raise ValueError("Extraction yielded no content")

    result["artifact_id"] = artifact_id
    result["extracted_at"] = datetime.now().isoformat()
    result["source_format"] = ext
    result["full_text"] = _build_full_text(result["sections"])
    result["word_count"] = len(result["full_text"].split())

    extracted_path = config.field_dir / "extracted" / f"{artifact_id}.json"
    _atomic_json(extracted_path, result)
    result["extracted_path"] = str(extracted_path)

    _log.info(
        "extraction_complete",
        artifact_id=artifact_id,
        ext=ext,
        sections=len(result["sections"]),
        word_count=result["word_count"],
        notes=len(result["extraction_notes"]),
    )
    return result


# ── Format extractors ────────────────────────────────────────────────────────

def _extract_docx(path: Path) -> dict:
    from docx import Document
    from docx.oxml.ns import qn

    sections = []
    tables = []
    notes = []
    title_extracted = None

    try:
        doc = Document(str(path))
    except Exception as exc:
        notes.append(f"Failed to open docx: {exc}")
        return _empty(notes)

    current_section_idx = 0

    for elem in doc.element.body:
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

        if tag == "p":
            try:
                from docx.text.paragraph import Paragraph
                para = Paragraph(elem, doc)
                text = para.text.strip()
                if not text:
                    continue
                style_name = para.style.name if para.style else ""
                level = _heading_level_from_style(style_name)
                if level is not None:
                    if level == 1 and title_extracted is None:
                        title_extracted = text
                    sections.append({"heading": text, "level": level, "content": ""})
                    current_section_idx = len(sections) - 1
                else:
                    if sections:
                        sep = "\n" if sections[-1]["content"] else ""
                        sections[-1]["content"] += sep + text
                    else:
                        sections.append({"heading": "", "level": 0, "content": text})
                        current_section_idx = 0
            except Exception as exc:
                notes.append(f"Paragraph parse error: {exc}")

        elif tag == "tbl":
            try:
                from docx.table import Table
                tbl = Table(elem, doc)
                rows = []
                for row in tbl.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                tables.append({"section_index": current_section_idx, "rows": rows})
            except Exception as exc:
                notes.append(f"Table parse error: {exc}")

    return {
        "title_extracted": title_extracted,
        "sections": sections,
        "tables": tables,
        "extraction_notes": notes,
    }


def _extract_xlsx(path: Path) -> dict:
    import openpyxl

    sections = []
    tables = []
    notes = []
    title_extracted = None

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        notes.append(f"Failed to open xlsx: {exc}")
        return _empty(notes)

    for sheet_idx, sheet_name in enumerate(wb.sheetnames):
        try:
            ws = wb[sheet_name]
            if sheet_idx == 0:
                title_extracted = sheet_name

            rows = []
            for row in ws.iter_rows(values_only=True):
                str_row = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in str_row):
                    rows.append(str_row)

            # Separate first row as headers if it contains string-only values
            # (i.e. looks like a header row rather than data)
            headers = None
            data_rows = rows
            if rows:
                first = rows[0]
                if first and all(
                    isinstance(cell, str) and not cell.replace(".", "").replace("-", "").isdigit()
                    for cell in first if cell.strip()
                ):
                    headers = first
                    data_rows = rows[1:]

            content_lines = ["\t".join(r) for r in rows]
            sections.append({
                "heading": sheet_name,
                "level": 1,
                "content": "\n".join(content_lines),
            })
            tables.append({
                "section_index": len(sections) - 1,
                "headers": headers,
                "rows": data_rows,
            })
        except Exception as exc:
            notes.append(f"Sheet {sheet_name!r} parse error: {exc}")

    wb.close()
    return {
        "title_extracted": title_extracted,
        "sections": sections,
        "tables": tables,
        "extraction_notes": notes,
    }


def _extract_pdf(path: Path) -> dict:
    import pdfplumber

    sections = []
    tables = []
    notes = []
    title_extracted = None

    try:
        pdf = pdfplumber.open(str(path))
    except Exception as exc:
        notes.append(f"Failed to open pdf: {exc}")
        return _empty(notes)

    try:
        # First pass: collect all chars to find median font size for heading detection
        all_sizes = []
        for page in pdf.pages:
            try:
                chars = page.chars
                all_sizes.extend(c.get("size", 0) for c in chars if c.get("size"))
            except Exception:
                pass

        median_size = sorted(all_sizes)[len(all_sizes) // 2] if all_sizes else 0
        heading_threshold = median_size * 1.25 if median_size else 0

        for page_num, page in enumerate(pdf.pages, start=1):
            try:
                text = (page.extract_text() or "").strip()
                if not text:
                    notes.append(f"Page {page_num}: little or no text extracted (image-heavy?)")

                # Try font-size-based heading detection
                if heading_threshold > 0 and page.chars:
                    lines = _pdf_lines_with_sizes(page)
                    page_sections = _pdf_sections_from_lines(lines, heading_threshold)
                    for i, sec in enumerate(page_sections):
                        if sec["heading"] and page_num == 1 and title_extracted is None:
                            title_extracted = sec["heading"]
                        sections.append(sec)
                else:
                    if page_num == 1 and text:
                        first_line = text.splitlines()[0].strip()
                        if first_line:
                            title_extracted = first_line
                    sections.append({"heading": f"Page {page_num}", "level": 1, "content": text})

                # Extract tables
                try:
                    for tbl in page.extract_tables() or []:
                        if tbl:
                            rows = [[str(c) if c is not None else "" for c in row] for row in tbl]
                            tables.append({"section_index": len(sections) - 1, "rows": rows})
                except Exception as exc:
                    notes.append(f"Page {page_num} table extract error: {exc}")

            except Exception as exc:
                notes.append(f"Page {page_num} parse error: {exc}")
    finally:
        pdf.close()

    # If pdfplumber yielded very little text, try PyMuPDF as fallback.
    # Browser-print PDFs from news sites often have body text invisible to pdfplumber.
    total_words = sum(len(s.get("content", "").split()) for s in sections)
    page_count = max(1, len(sections))
    if total_words < 150 or (total_words / page_count) < 30:
        fitz_result = _extract_pdf_fitz(path)
        fitz_words = sum(len(s.get("content", "").split()) for s in fitz_result.get("sections", []))
        if fitz_words > total_words:
            fitz_result["extraction_notes"] = (
                [f"pdfplumber yielded {total_words} words; PyMuPDF fallback yielded {fitz_words} words — using PyMuPDF"]
                + fitz_result.get("extraction_notes", [])
            )
            return fitz_result
        notes.append(f"pdfplumber sparse ({total_words} words); PyMuPDF fallback also sparse ({fitz_words} words) — PDF may be image-based")

    return {
        "title_extracted": title_extracted,
        "sections": sections,
        "tables": tables,
        "extraction_notes": notes,
    }


def _extract_pdf_fitz(path: Path) -> dict:
    """PyMuPDF-based PDF extraction. Better than pdfplumber for browser-print PDFs."""
    import fitz  # PyMuPDF

    sections = []
    notes = []
    title_extracted = None

    try:
        doc = fitz.open(str(path))
    except Exception as exc:
        notes.append(f"PyMuPDF failed to open: {exc}")
        return _empty(notes)

    try:
        # Collect all font sizes across all pages to find median for heading detection
        all_sizes = []
        for page in doc:
            try:
                blocks = page.get_text("dict").get("blocks", [])
                for b in blocks:
                    for line in b.get("lines", []):
                        for span in line.get("spans", []):
                            sz = span.get("size", 0)
                            if sz > 0:
                                all_sizes.append(sz)
            except Exception:
                pass

        median_size = sorted(all_sizes)[len(all_sizes) // 2] if all_sizes else 0
        heading_threshold = median_size * 1.25 if median_size else 0

        current_heading = ""
        current_level = 1
        current_parts: list[str] = []

        def flush_section():
            nonlocal current_heading, current_level, current_parts
            content = "\n".join(current_parts).strip()
            if current_heading or content:
                sections.append({"heading": current_heading, "level": current_level, "content": content})
            current_heading = ""
            current_level = 1
            current_parts = []

        for page_num, page in enumerate(doc, start=1):
            try:
                blocks = page.get_text("dict").get("blocks", [])
                for b in blocks:
                    if b.get("type") != 0:  # skip image blocks
                        continue
                    for line in b.get("lines", []):
                        line_text = "".join(s.get("text", "") for s in line.get("spans", [])).strip()
                        if not line_text:
                            continue
                        max_size = max((s.get("size", 0) for s in line.get("spans", [])), default=0)
                        if heading_threshold > 0 and max_size >= heading_threshold:
                            flush_section()
                            current_heading = line_text
                            if page_num == 1 and title_extracted is None:
                                title_extracted = line_text
                        else:
                            current_parts.append(line_text)
            except Exception as exc:
                notes.append(f"PyMuPDF page {page_num} error: {exc}")

        flush_section()

        # If heading detection produced nothing useful, fall back to plain page text
        if not sections:
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text().strip()
                if text:
                    if page_num == 1 and title_extracted is None:
                        title_extracted = text.splitlines()[0].strip()
                    sections.append({"heading": f"Page {page_num}", "level": 1, "content": text})

    finally:
        doc.close()

    return {
        "title_extracted": title_extracted,
        "sections": sections,
        "tables": [],
        "extraction_notes": notes,
    }


def _pdf_lines_with_sizes(page) -> list[dict]:
    """Group PDF chars into lines, tracking max font size per line."""
    lines: list[dict] = []
    current_line: list = []
    current_top = None
    TOL = 3  # px tolerance for same line

    for char in sorted(page.chars, key=lambda c: (round(c["top"] / TOL), c["x0"])):
        top = round(char["top"] / TOL) * TOL
        if current_top is None:
            current_top = top
        if abs(top - current_top) > TOL:
            if current_line:
                text = "".join(c.get("text", "") for c in current_line).strip()
                max_size = max((c.get("size", 0) for c in current_line), default=0)
                if text:
                    lines.append({"text": text, "size": max_size})
            current_line = [char]
            current_top = top
        else:
            current_line.append(char)

    if current_line:
        text = "".join(c.get("text", "") for c in current_line).strip()
        max_size = max((c.get("size", 0) for c in current_line), default=0)
        if text:
            lines.append({"text": text, "size": max_size})
    return lines


def _pdf_sections_from_lines(lines: list[dict], heading_threshold: float) -> list[dict]:
    """Convert lines with sizes into sections, treating large text as headings."""
    sections: list[dict] = []
    current_heading = ""
    current_level = 1
    current_parts: list[str] = []

    def flush():
        if current_heading or current_parts:
            sections.append({
                "heading": current_heading,
                "level": current_level,
                "content": "\n".join(current_parts).strip(),
            })

    for line in lines:
        if line["size"] >= heading_threshold:
            flush()
            current_heading = line["text"]
            current_level = 1
            current_parts = []
        else:
            current_parts.append(line["text"])

    flush()
    return sections or [{"heading": "", "level": 0, "content": "\n".join(l["text"] for l in lines)}]


def _extract_pptx(path: Path) -> dict:
    from pptx import Presentation
    from pptx.util import Pt

    sections = []
    notes = []
    title_extracted = None

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        notes.append(f"Failed to open pptx: {exc}")
        return _empty(notes)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        try:
            slide_title = None
            content_parts = []

            for shape in slide.shapes:
                try:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    # First text box on first slide with content = title
                    if slide_title is None:
                        slide_title = text
                        if slide_idx == 1:
                            title_extracted = text
                    else:
                        content_parts.append(text)
                except Exception as exc:
                    notes.append(f"Slide {slide_idx} shape error: {exc}")

            sections.append({
                "heading": slide_title or f"Slide {slide_idx}",
                "level": 1,
                "content": "\n".join(content_parts),
            })
        except Exception as exc:
            notes.append(f"Slide {slide_idx} parse error: {exc}")

    return {
        "title_extracted": title_extracted,
        "sections": sections,
        "tables": [],
        "extraction_notes": notes,
    }


def _extract_markdown(path: Path) -> dict:
    sections = []
    notes = []
    title_extracted = None

    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        notes.append(f"Failed to read file: {exc}")
        return _empty(notes)

    current_heading = ""
    current_level = 0
    current_lines = []

    def flush():
        if current_heading or current_lines:
            sections.append({
                "heading": current_heading,
                "level": current_level,
                "content": "\n".join(current_lines).strip(),
            })

    for line in text.splitlines():
        m = re.match(r"^(#{1,6})\s+(.*)", line)
        if m:
            flush()
            current_level = len(m.group(1))
            current_heading = m.group(2).strip()
            current_lines = []
            if current_level == 1 and title_extracted is None:
                title_extracted = current_heading
        else:
            current_lines.append(line)

    flush()

    if not sections:
        sections.append({"heading": "", "level": 0, "content": text.strip()})

    return {
        "title_extracted": title_extracted,
        "sections": sections,
        "tables": [],
        "extraction_notes": notes,
    }


def _extract_txt(path: Path) -> dict:
    notes = []
    title_extracted = None

    try:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
    except Exception as exc:
        notes.append(f"Failed to read file: {exc}")
        return _empty(notes)

    first_line = text.splitlines()[0].strip() if text else ""
    if first_line:
        title_extracted = first_line[:120]

    return {
        "title_extracted": title_extracted,
        "sections": [{"heading": "", "level": 0, "content": text}],
        "tables": [],
        "extraction_notes": notes,
    }


# ── Helpers ──────────────────────────────────────────────────────────────────

def _heading_level_from_style(style_name: str) -> int | None:
    """Return heading level (1–6) from a Word style name, or None if not a heading."""
    m = re.match(r"[Hh]eading\s+(\d)", style_name)
    if m:
        return int(m.group(1))
    return None


def _build_full_text(sections: list[dict]) -> str:
    parts = []
    for s in sections:
        if s.get("heading"):
            parts.append(s["heading"])
        if s.get("content"):
            parts.append(s["content"])
    return "\n\n".join(parts)


def _atomic_json(path: Path, data: dict) -> None:
    """Write JSON atomically: write to .tmp, fsync, then os.replace."""
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    try:
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, default=str)
            f.flush()
            os.fsync(f.fileno())
        os.replace(tmp, path)
    except Exception:
        try:
            tmp.unlink(missing_ok=True)
        except Exception:
            pass
        raise


def _empty(notes: list) -> dict:
    return {
        "title_extracted": None,
        "sections": [],
        "tables": [],
        "extraction_notes": notes,
    }
